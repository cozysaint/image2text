import os
from pptx import Presentation
import tempfile
import pytesseract
from PIL import Image

def extract_text_from_image(image_path, lang='kor'):
    # 이미지에서 텍스트 추출
    img = Image.open(image_path)
    text = pytesseract.image_to_string(img, lang=lang)
    return text

def replace_images_with_text(pptx_file):
    prs = Presentation(pptx_file)

    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.shape_type == 13:  # 이미지 형태 (13은 사각형을 의미)
                with tempfile.NamedTemporaryFile(suffix=".jpg", delete=False) as tmp_file:
                    img_path = tmp_file.name  # 임시 이미지 저장 경로
                    with open(img_path, "wb") as f:
                        f.write(shape.image.blob)  # 이미지 데이터를 파일로 저장
                text = extract_text_from_image(img_path, lang='kor')
                os.remove(img_path)  # 임시 이미지 파일 삭제

                # 텍스트 상자로 변환
                left, top, width, height = shape.left, shape.top, shape.width, shape.height
                txBox = slide.shapes.add_textbox(left, top, width, height)
                tf = txBox.text_frame
                p = tf.add_paragraph()
                p.text = text

                # 원래 이미지 삭제
                slide.shapes._spTree.remove(shape._element)

    prs.save("output.pptx")  # 결과 저장

if __name__ == "__main__":
    pptx_file = "/Users/sungichul/Desktop/sungichul/coding/ppt/test.pptx"  # 입력 PPT 파일 경로
    replace_images_with_text(pptx_file)

